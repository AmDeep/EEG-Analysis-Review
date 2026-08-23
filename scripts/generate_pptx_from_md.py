"""
scripts/generate_pptx_from_md.py

Enhanced generator: converts a simple markdown slide deck (docs/pipeline_slides.md)
into a styled PowerPoint (docs/pipeline_slides.pptx) with colors, fonts, and optional
local images embedded. Designed for the EEG-Analysis-Review repository slide deck.

Usage:
    python scripts/generate_pptx_from_md.py docs/pipeline_slides.md docs/pipeline_slides.pptx

Dependencies:
    pip install python-pptx Pillow

Behavior:
- Treats '---' as slide separator.
- First H1 (# ...) becomes title slide title; subsequent H2/lines become subtitle or bullets.
- Recognizes bullet lines starting with '- ' or '* '.
- If a line references an image path (endswith .png/.jpg/.jpeg), the script will attempt to
  embed that image (path relative to repo root).
- Adds a consistent color scheme and fonts; places speaker notes with the full slide text.

This script intentionally keeps parsing simple but improves formatting compared to the
original lightweight version.
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from PIL import Image

# Styling constants
TITLE_FONT = 'Calibri'
BODY_FONT = 'Calibri'
TITLE_FONT_SIZE = Pt(40)
SUBTITLE_FONT_SIZE = Pt(20)
BODY_FONT_SIZE = Pt(18)
ACCENT_COLOR = RGBColor(10, 84, 126)       # deep teal/navy
ACCENT_COLOR_DARK = RGBColor(6, 47, 70)
BG_COLOR = RGBColor(245, 247, 250)        # very light gray
TITLE_COLOR = RGBColor(255, 255, 255)
BODY_COLOR = RGBColor(20, 20, 20)
HIGHLIGHT = RGBColor(0, 123, 167)         # bright teal for bullets

IMAGE_MAX_WIDTH = Inches(9)
IMAGE_MAX_HEIGHT = Inches(4.5)


def md_to_slides(md_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        lines = [l.rstrip('\n') for l in f]

    slides = []
    current = []
    for line in lines:
        if line.strip() == '---':
            if current:
                slides.append(current)
                current = []
        else:
            current.append(line)
    if current:
        slides.append(current)
    return slides


def is_image_line(line):
    low = line.lower()
    return low.endswith('.png') or low.endswith('.jpg') or low.endswith('.jpeg') or low.startswith('![')


def extract_image_path(line):
    # supports both plain paths and Markdown image syntax ![alt](path)
    line = line.strip()
    if line.startswith('!['):
        # find the parentheses
        start = line.find('(')
        end = line.rfind(')')
        if start != -1 and end != -1 and end > start:
            return line[start+1:end]
        return None
    return line


def add_title_slide(prs, title_text, subtitle_text=None):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = ACCENT_COLOR

    # Title box
    left = Inches(0.6)
    top = Inches(1.2)
    width = Inches(8.0)
    height = Inches(1.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = TITLE_FONT
    p.font.size = TITLE_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    if subtitle_text:
        st_top = Inches(3.0)
        st_box = slide.shapes.add_textbox(left, st_top, width, Inches(1.0))
        st_tf = st_box.text_frame
        st_p = st_tf.paragraphs[0]
        st_p.text = subtitle_text
        st_p.font.name = BODY_FONT
        st_p.font.size = SUBTITLE_FONT_SIZE
        st_p.font.bold = False
        st_p.font.color.rgb = TITLE_COLOR

    return slide


def add_bullet_slide(prs, title, bullets, images=None):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # Title
    left = Inches(0.4)
    top = Inches(0.3)
    width = Inches(9.2)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_tf = title_box.text_frame
    t = title_tf.paragraphs[0]
    t.text = title
    t.font.name = TITLE_FONT
    t.font.size = Pt(28)
    t.font.bold = True
    t.font.color.rgb = ACCENT_COLOR_DARK

    # Content box
    c_top = Inches(1.2)
    c_height = Inches(3.6)
    content_box = slide.shapes.add_textbox(left, c_top, width, c_height)
    content_tf = content_box.text_frame
    content_tf.word_wrap = True

    first = True
    for b in bullets:
        if first:
            p = content_tf.paragraphs[0]
            p.text = b
            first = False
        else:
            p = content_tf.add_paragraph()
            p.text = b
        p.level = 0
        p.font.name = BODY_FONT
        p.font.size = BODY_FONT_SIZE
        p.font.color.rgb = BODY_COLOR

    # If images provided, place one on the lower area
    if images:
        # Try to place the first image on the right side
        img_path = images[0]
        if os.path.exists(img_path):
            try:
                with Image.open(img_path) as im:
                    width_px, height_px = im.size
                pic_left = Inches(0.6)
                pic_top = Inches(4.6)
                pic_width = IMAGE_MAX_WIDTH
                pic_height = IMAGE_MAX_HEIGHT
                slide.shapes.add_picture(img_path, pic_left, pic_top, width=pic_width, height=pic_height)
            except Exception:
                pass

    # Add speaker notes with full bullets for reference
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = title + '\n\n' + '\n'.join(bullets)

    return slide


def prepare_slide_content(block):
    title = None
    subtitle = None
    bullets = []
    images = []

    for line in block:
        s = line.strip()
        if not s:
            continue
        if s.startswith('# '):
            if title is None:
                title = s.lstrip('#').strip()
            else:
                bullets.append(s.lstrip('#').strip())
        elif s.startswith('## '):
            # treat as subtitle if title empty
            if title is None:
                title = s.lstrip('#').strip()
            else:
                bullets.append(s.lstrip('#').strip())
        elif s.startswith('- ') or s.startswith('* '):
            bullets.append(s[2:].strip())
        elif is_image_line(s):
            img = extract_image_path(s)
            if img:
                images.append(img)
        else:
            # fallback: treat as bullet
            bullets.append(s)

    return title, subtitle, bullets, images


def make_pptx(md_path, pptx_path):
    slides_blocks = md_to_slides(md_path)
    prs = Presentation()

    # metadata: set slide size to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for i, block in enumerate(slides_blocks):
        title, subtitle, bullets, images = prepare_slide_content(block)
        # Use first block as title slide
        if i == 0:
            if not title:
                title = 'Title'
            # Combine remaining bullets into subtitle if present
            subtitle_text = None
            if bullets:
                subtitle_text = bullets[0]
            add_title_slide(prs, title, subtitle_text)
        else:
            slide_title = title if title else f'Slide {i+1}'
            # If no bullets, try to use block lines as bullets
            if not bullets and not images:
                bullets = [l.strip() for l in block if l.strip()]
            add_bullet_slide(prs, slide_title, bullets, images)

    prs.save(pptx_path)
    print(f'Wrote {pptx_path}')


if __name__ == '__main__':
    if len(sys.argv) != 3:
        print('Usage: python scripts/generate_pptx_from_md.py input.md output.pptx')
        sys.exit(1)
    md_path = sys.argv[1]
    pptx_path = sys.argv[2]
    make_pptx(md_path, pptx_path)
