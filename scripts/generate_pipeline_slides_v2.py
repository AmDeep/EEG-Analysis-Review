"""
generate_pipeline_slides_v2.py

Creates docs/pipeline_slides_v2.pptx from the design spec in docs/pipeline_slides_v2.md.

Usage (run locally):
  1. Create a Python venv (optional): python -m venv .venv && source .venv/bin/activate
  2. Install dependencies: pip install python-pptx Pillow
  3. Run: python scripts/generate_pipeline_slides_v2.py

The script writes docs/pipeline_slides_v2.pptx. Optional: provide a logo file at assets/logo.png to include on the title slide.

Note: This environment cannot run arbitrary binaries; commit this script and run it locally or in your CI to produce the PPTX.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

OUT_PATH = "docs/pipeline_slides_v2.pptx"
LOGO_PATH = "assets/logo.png"  # optional

# Simple helper to add a title slide

def add_title_slide(prs, title, subtitle, palette):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(44)
    subtitle_tf = slide.placeholders[1].text_frame
    subtitle_tf.text = subtitle
    subtitle_tf.paragraphs[0].font.size = Pt(18)
    # add logo if present
    if os.path.exists(LOGO_PATH):
        try:
            slide.shapes.add_picture(LOGO_PATH, prs.slide_width - Inches(1.5), Inches(0.3), width=Inches(1.2))
        except Exception:
            pass
    return slide


def add_bulleted_slide(prs, title, bullets, notes=None, palette=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i>0 else body.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(20)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_two_column_slide(prs, title, left_lines, right_lines, notes=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    # remove default body
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.2), Inches(4.5))
    lf = left.text_frame
    for i, l in enumerate(left_lines):
        p = lf.add_paragraph() if i>0 else lf.paragraphs[0]
        p.text = l
        p.font.size = Pt(18)
    right = slide.shapes.add_textbox(Inches(5.1), Inches(1.6), Inches(4.2), Inches(4.5))
    rf = right.text_frame
    for i, r in enumerate(right_lines):
        p = rf.add_paragraph() if i>0 else rf.paragraphs[0]
        p.text = r
        p.font.size = Pt(18)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def main():
    # color palette defaults (RGB tuples)
    palette = {
        'primary': RGBColor(0x3F,0x51,0xB5),
        'secondary': RGBColor(0x00,0x96,0x88),
        'accent': RGBColor(0xFF,0x6F,0x61)
    }

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
                    "EEG Analysis Pipeline — Feature Overview & Design (v2)",
                    "Reworked visuals, clarified features, integrated Hrishi notes\nRepo: AmDeep/EEG-Analysis-Review",
                    palette)

    # Slide 2: Executive summary
    add_bulleted_slide(prs,
                       "Executive summary / Key takeaways",
                       [
                           "What this pipeline extracts: EPs (N2/P2), time-frequency, connectivity, entropy/complexity, metadata.",
                           "Why it matters: N2-P2 and gamma are validated pain biomarkers; connectivity and entropy add mechanistic insight.",
                           "Next steps: fill missing ERD, standardize gamma band, restrict connectivity to compatible montages."
                       ],
                       notes="Speaker notes: High-level context and immediate actions."
                       )

    # Slide 3: Data provenance & epoching
    add_bulleted_slide(prs,
                       "Epochs & Timing",
                       [
                           "Pre-epoched ~3s windows around each laser pulse (dataset authors).",
                           "Epoch column indexes pre-segmented trials.",
                           "Inter-epoch timing varies (~3–8s jitter) plus fixation/rating extends gap."
                       ],
                       notes="Hrishi Q1/Q2: epochs are 3s bursts; ITIs jitter 3–8s plus extra fixation/rating time."
                       )

    # Slide 4: Trial interaction
    add_bulleted_slide(prs,
                       "Trial interactions & translational limits",
                       [
                           "Repeated stimulation → peripheral sensitization + neural habituation (N2–P2 amplitude shrink).",
                           "Some sub-experiments rotate stimulation site to reduce skin sensitization.",
                           "Not equivalent to chronic pain cohorts: single-session healthy volunteers; mechanistic overlap only."
                       ],
                       notes="Hrishi Q3/Q4: describe sensitization, habituation, and limits vs chronic pain."
                       )

    # Slide 5: Feature categories overview
    add_two_column_slide(prs,
                         "Feature categories",
                         [
                             "Evoked potentials: n2_amp, n2_lat, p2_amp, p2_lat, n2p2_amp",
                             "Oscillatory & spectral: gamma_power, alpha_erd_pct, beta_erd_pct, psd_*"
                         ],
                         [
                             "Connectivity: plv_Fz-Cz, plv_Cz-Pz, plv_C3-C4, plv_FCz-CPz",
                             "Complexity/entropy: perm_entropy, spectral_entropy, sample_entropy, higuchi_fd, dfa, hjorth_*"
                         ],
                         notes="Group features by measurement type."
                         )

    # Slides 6-12: Feature group slides (compact)
    add_bulleted_slide(prs,
                       "Evoked potentials (N2/P2)",
                       [
                           "n2_amp, n2_lat — N2 negative deflection (~150–350ms).",
                           "p2_amp, p2_lat — P2 positive deflection (~300–550ms).",
                           "n2p2_amp — peak-to-peak (p2_amp - n2_amp): strongest pain biomarker."
                       ],
                       notes="Show waveform example in presenter view; discuss units (µV) and latencies."
                       )

    add_bulleted_slide(prs,
                       "Gamma & high-frequency features",
                       [
                           "gamma_power (per gamma_band_hz): reflects pain salience; vulnerable to muscle artifact.",
                           "psd_gamma: raw gamma power via Welch's method (complementary)."
                       ],
                       notes="Caution about EMG contamination; consider visual QC and EMG regressors."
                       )

    add_bulleted_slide(prs,
                       "Alpha/Beta ERD & PSD",
                       [
                           "alpha_erd_pct: percent reduction vs pre-stim baseline (currently 100% missing).",
                           "beta_erd_pct: percent reduction in beta band; psd_alpha/psd_beta are absolute powers."
                       ],
                       notes="ERD% vs PSD difference and options to compute/impute."
                       )

    add_bulleted_slide(prs,
                       "Low-frequency PSD features",
                       [
                           "psd_delta (1–4Hz): drowsiness marker; sometimes secondary pain marker.",
                           "psd_theta (4–8Hz): increases with pain intensity in acute and chronic settings."
                       ],
                       notes="Discuss Sarnthein reference for theta in chronic pain."
                       )

    add_bulleted_slide(prs,
                       "Connectivity features (PLV)",
                       [
                           "plv_Fz-Cz, plv_Cz-Pz, plv_C3-C4, plv_FCz-CPz (FCz-CPz ~85% missing).",
                           "Phase-locking measures capture interregional timing consistency; often informative vs power."
                       ],
                       notes="Recommend restricting PLV analyses to compatible montages; document missingness."
                       )

    add_bulleted_slide(prs,
                       "Complexity & entropy measures",
                       [
                           "perm_entropy, spectral_entropy, sample_entropy: signal unpredictability metrics.",
                           "higuchi_fd, dfa: fractal & long-range temporal correlation measures; hjorth_mobility/complexity: simple proxies for frequency & irregularity."
                       ],
                       notes="These features summarize nonstationary/irregular aspects of the EEG."
                       )

    add_bulleted_slide(prs,
                       "Metadata & labels",
                       [
                           "sfreq: sampling frequency (Hz).",
                           "gamma_band_hz: metadata specifying gamma sub-band (e.g., '30-80').",
                           "rating (0–10 NRS): subjective pain label; laser_power (Joules): objective stimulus intensity."
                       ],
                       notes="rating is the main ground-truth target for models."
                       )

    # Slide 13: Missingness & data-quality
    add_bulleted_slide(prs,
                       "Missingness & Data Quality",
                       [
                           "alpha_erd_pct: 100% missing in extracted table.",
                           "plv_FCz-CPz: ~85% missing due to montage differences.",
                           "Options: impute, compute from raw if available, or restrict analyses."
                       ],
                       notes="Recommend prioritizing features with low missingness for baseline models."
                       )

    # Slide 14: Preprocessing & confounds
    add_bulleted_slide(prs,
                       "Preprocessing & Confounds to Control",
                       [
                           "Filters: notch line noise, bandpass 0.5–120Hz; baseline correction; ICA for artifact removal.",
                           "Confound regressors: trial number (habituation), stimulation site, subject intercept, sfreq differences, laser_power."
                       ],
                       notes="Include trial-order regressors to account for habituation/sensitization effects."
                       )

    # Slide 15: Example analysis pipeline
    add_bulleted_slide(prs,
                       "Analysis pipeline (example)",
                       [
                           "Load → preprocess → epoch QC → feature extraction → feature QC/missingness handling → model (rating) → explainability.",
                           "Add explainability: feature importance, partial dependence plots, per-sub-experiment reporting."
                       ],
                       notes="Link to notebook in repo for runnable code."
                       )

    # Slide 16: Example epoch timeline
    add_bulleted_slide(prs,
                       "Epoch timeline (worked example)",
                       [
                           "Pre-stim baseline: -500–0ms; stimulus onset: 0ms; N2 window: 150–350ms; P2 window: 300–550ms.",
                           "sfreq example: 1000Hz → 1ms samples; latency units reported in seconds in table."
                       ],
                       notes="Use this to explain latency columns and windowing choices."
                       )

    # Slide 17: Appendix A: Full feature definitions (Hrishi)
    appendix_a = [
        "sfreq — The sampling frequency (Hz) the EEG was recorded at for that trial/dataset (e.g., 1000Hz or 1024Hz). It's metadata rather than a physiological feature — it just tells you the time resolution of the recording, and different sub-experiments in the collection use slightly different values depending on which EEG system recorded them.",
        "n2_amp — The amplitude (µV) of the N2 wave, a negative voltage deflection ~150–350ms after a painful stimulus. Reflects how strongly the sensory cortex reacted, and is one half of the classic laser-evoked potential pair.",
        "n2_lat — The timing (seconds) of that N2 peak. Faster latency generally reflects quicker A-delta fiber conduction speed.",
        "p2_amp — The amplitude of the P2 wave, a positive deflection ~300–550ms post-stimulus. The most consistently pain-intensity-correlated single ERP component in the literature.",
        "p2_lat — The timing of the P2 peak, used alongside p2_amp to describe the full late cortical response shape.",
        "n2p2_amp — The peak-to-peak amplitude between N2 and P2 (p2_amp minus n2_amp). The single most commonly cited pain-intensity biomarker in laser-EEG research.",
        "gamma_power — Power in the gamma band after the stimulus, computed over the specific frequency range given in gamma_band_hz. Thought to reflect moment-to-moment pain salience, though vulnerable to muscle-artifact contamination.",
        "gamma_band_hz — The exact frequency range (e.g., \"30-80\") used to compute gamma power for that row. Not a feature itself — it's metadata documenting which gamma sub-band was used, since \"gamma\" isn't strictly standardized across studies.",
        "alpha_erd_pct — Percent reduction in alpha power after the stimulus vs. a pre-stimulus baseline. Bigger drops mean the brain is more actively engaged; currently 100% missing in your extracted data.",
        "beta_erd_pct — Same idea as alpha ERD, but for the beta band — linked to disengaging the current motor state as the brain prepares a protective response.",
        "psd_delta — Delta-band (1–4Hz) power via Welch's method. More associated with drowsiness generally, but appears as a secondary marker in some chronic-pain resting-state work.",
        "psd_theta — Theta-band (4–8Hz) power. Rises with pain intensity in acute settings, and is one of the strongest chronic-pain resting-state markers (Sarnthein et al.).",
        "psd_alpha — Alpha-band raw power via Welch PSD (distinct from the ERD percentage — this is absolute power, not percent-change). Useful both acutely and at rest.",
        "psd_beta — Beta-band raw power via Welch PSD. Raw-power counterpart to beta_erd_pct.",
        "psd_gamma — Gamma-band raw power via Welch PSD, the raw-power counterpart to gamma power. Frequently one of the most informative features in pain-classification ML work.",
        "plv_Fz-Cz — Phase-locking value between frontal (Fz) and central (Cz) electrodes — how consistently their oscillation timing stays aligned. Phase-connectivity features have outperformed simple power features in recent pain-classification work.",
        "plv_Cz-Pz — Same phase-locking concept between central (Cz) and parietal (Pz) electrodes, capturing front-back connectivity.",
        "plv_C3-C4 — Phase-locking between left (C3) and right (C4) central electrodes — tests interhemispheric connectivity and laterality, since pain responses are typically stronger contralateral to the stimulated hand.",
        "plv_FCz-CPz — Phase-locking between FCz and CPz, electrodes just anterior/posterior to the vertex (Cz) — another connectivity angle centered on the primary sensorimotor region. Currently 85% missing since not every montage includes both electrodes.",
        "perm_entropy — Permutation entropy — how unpredictable the signal's ordering pattern is. Higher values mean a more irregular, less repetitive signal.",
        "spectral_entropy — Entropy computed on the frequency spectrum rather than the raw waveform — how spread out vs. concentrated the power is across frequencies.",
        "sample_entropy — Measures how likely similar patterns in the signal are to repeat. Tied to trial-to-trial variability, which has independently predicted pain perception in prior work.",
        "higuchi_fd — Higuchi fractal dimension — a geometric complexity measure of the waveform (how much fine detail/roughness it has across scales).",
        "dfa — Detrended fluctuation analysis exponent — measures long-range temporal correlations, i.e., whether fluctuations early in the signal relate statistically to fluctuations much later.",
        "hjorth_mobility — One of the three Hjorth parameters; approximates the signal's mean frequency, computed cheaply from variance ratios without a full frequency transform.",
        "hjorth_complexity — The second Hjorth parameter; measures how much the signal's frequency content changes over time — a proxy for how irregular vs. sine-wave-like the waveform is.",
        "rating — The subjective pain rating (typically 0–10 NRS) the participant gave for that specific trial. This is your ground-truth label for \"how much did it hurt.\"",
        "laser_power — The physical stimulus intensity (in Joules) delivered on that trial. This is the objective stimulus-strength label, used in the notebook as the comparison target against rating"
    ]
    add_bulleted_slide(prs, "Appendix A: Feature definitions (Hrishi)", appendix_a, notes="Full definitions included verbatim from Hrishi's notes.")

    # Slide 18: Appendix B: Q&A (Hrishi)
    appendix_b = [
        "1. Is epoch self-extracted, or 3-second bursts? 3-second bursts — the source data comes pre-epoched into ~3-second windows around each laser pulse by the original dataset authors. Epoch column indexes through those already-segmented trials.",
        "2. Time between epochs? Varies by sub-experiment, roughly 3–8 seconds of jittered inter-trial interval, plus a few extra seconds for fixation and rating in between. So the real gap between stimulus onsets is longer than the ITI number alone suggests.",
        "3. Interaction between trials? Yes — repeated stimulation causes both physical skin sensitization/fatigue and neural habituation (shrinking N2-P2 amplitude over the session), even when the person's reported pain doesn't change. Therefore, some sub-experiments rotate the stimulation site instead of hitting the same spot repeatedly.",
        "4. Does this make it resemble chronic pain? Partially — repeated-trial sensitization is mechanistically related to central sensitization, one driver of acute-to-chronic transition. But it's still a single-session, healthy-volunteer paradigm, not a real stand-in for actual chronic pain data(not too much data to assume)"
    ]
    add_bulleted_slide(prs, "Appendix B: Q&A (Hrishi)", appendix_b, notes="Q&A from Hrishi. Key takeaway: pre-epoched 3s windows; ITIs jittered; habituation/sensitization present; not chronic data.")

    # Slide 19: Next steps
    add_bulleted_slide(prs,
                       "Next steps & action items",
                       [
                           "Add ERD extraction or compute and populate alpha_erd_pct column.",
                           "Standardize gamma_band_hz across sub-experiments and recompute gamma_power.",
                           "Restrict connectivity analyses to montages with full electrode set; document per-sub-experiment availability.",
                           "Run habituation checks (trial-number regressors) and add plots to notebook."
                       ],
                       notes="Actionable items and owner assignments."
                       )

    # Save PPTX
    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved {OUT_PATH}")


if __name__ == '__main__':
    main()
